from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


def build_pptx_report(title: str, sections: list[tuple[str, str]], output_path: Path) -> None:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Enterprise Multi-Agent Operations Copilot"

    bullet_layout = prs.slide_layouts[1]
    for heading, content in sections:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = heading
        body = slide.placeholders[1].text_frame
        body.clear()
        lines = [line.strip("- ") for line in content.splitlines() if line.strip()][:8] or ["(no content)"]
        body.text = lines[0]
        for line in lines[1:]:
            p = body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)

    prs.save(str(output_path))
